# -*- coding: utf-8 -*-
"""Gera Apresentacao_Nutricao_Diretoria.pptx — São Geraldo Service."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "Apresentacao_Nutricao_Diretoria.pptx"

# Marca São Geraldo Service
BLUE = RGBColor(0x1B, 0x4F, 0x9C)
NAVY = RGBColor(0x0D, 0x2B, 0x5C)
LIME = RGBColor(0xA8, 0xC5, 0x3A)
LIME_DARK = RGBColor(0x8F, 0xAF, 0x28)
INK = RGBColor(0x0F, 0x2A, 0x56)
MIST = RGBColor(0xF3, 0xF7, 0xFC)
SOFT = RGBColor(0xE7, 0xEE, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_LINE = RGBColor(0xC5, 0xD0, 0xE0)


def _set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide):
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), LIME)


def _add_footer(slide, page, total=13):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8.5), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, f"São Geraldo Service  ·  Nutrição Hospitalar  ·  {page}/{total}", size=10, color=GRAY)
    # thin line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.95), Inches(9.0), Emu(5000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_LINE
    line.line.fill.background()


def _title_block(slide, title, subtitle=None):
    _add_accent_bar(slide)
    _add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.05), NAVY)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, title, size=26, bold=True, color=WHITE, font="Calibri")
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.4))
        p2 = sub.text_frame.paragraphs[0]
        r2 = p2.add_run()
        _set_run(r2, subtitle, size=14, color=GRAY)


def _bullets(slide, items, left=0.55, top=1.75, width=12.0, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        _set_run(run, f"•  {item}", size=size, color=INK)


def _two_columns(slide, left_items, right_items, top=1.85):
    _bullets(slide, left_items, left=0.55, top=top, width=5.8, size=15)
    _bullets(slide, right_items, left=6.7, top=top, width=5.8, size=15)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    _add_rect(slide, Inches(0), Inches(5.9), Inches(13.333), Inches(1.6), BLUE)
    _add_rect(slide, Inches(0), Inches(5.9), Inches(13.333), Inches(0.12), LIME)

    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6))
    r = t1.text_frame.paragraphs[0].add_run()
    _set_run(r, "SÃO GERALDO SERVICE", size=20, bold=True, color=LIME)

    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.0))
    r2 = t2.text_frame.paragraphs[0].add_run()
    _set_run(r2, "Nutrição Hospitalar", size=40, bold=True, color=WHITE)

    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.6))
    r3 = t3.text_frame.paragraphs[0].add_run()
    _set_run(r3, "Apresentação à Diretoria — visão de valor, operação e governança", size=16, color=SOFT)

    t4 = slide.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.8))
    tf = t4.text_frame
    p = tf.paragraphs[0]
    r4 = p.add_run()
    _set_run(r4, "Plataforma MeuApp  ·  Módulo integrado  ·  Uso interno", size=13, color=WHITE)


def slide_objetivo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Objetivo / problema que resolve", "Por que investir no módulo")
    _bullets(
        slide,
        [
            "Unificar a produção diária de refeições em um mapa oficial por clínica e enfermaria.",
            "Eliminar planilhas soltas e perda de histórico ao “apagar” pacientes.",
            "Padronizar dietas, horários (com hora limite), cardápios e preços.",
            "Dar rastreabilidade: motivo de saída, usuário e data/hora da alteração.",
            "Apoiar faturamento, etiquetas e totalizações a partir da mesma fonte de verdade.",
        ],
    )
    _add_footer(slide, 2)


def slide_visao(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Visão geral do módulo", "Um sistema, vários papéis")
    _two_columns(
        slide,
        [
            "Operação: Mapa de Produção do dia",
            "Cadastros clínicos e nutricionais",
            "Cardápios e dietas líquidas",
            "Preços (Funcionário / Paciente / Acompanhante)",
        ],
        [
            "Estoque, produtos e fornecedores",
            "Tabela de nutrientes + importação FDC",
            "Etiquetas, U.M.A., totalização",
            "Faturamento e análise de custos",
        ],
        top=1.9,
    )
    note = slide.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(12), Inches(0.8))
    r = note.text_frame.paragraphs[0].add_run()
    _set_run(
        r,
        "Acesso: login da plataforma → hub Sistemas → Nutrição  |  Atalho técnico: iniciar_meuapp.bat",
        size=13,
        color=GRAY,
    )
    _add_footer(slide, 3)


def slide_mapa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Mapa de Produção", "Coração da operação do dia")
    _bullets(
        slide,
        [
            "Filtros por clínica e enfermaria — a grade só carrega após o filtro (evita confusão e sobrecarga).",
            "Navegação dia a dia com persistência automática das linhas ativas.",
            "Flags de refeição: Desjejum, Colação, Almoço, Merenda, Jantar e Ceia.",
            "Campos complementares: extras, suplementos, enteral, fórmula infantil, LVE, observação de etiqueta.",
            "Substituições de cardápio com justificativa; avisos quando alta aparece em datas futuras.",
            "Inserção explícita de pacientes; exclusão somente com motivo (alta, óbito ou transferência).",
        ],
        size=15,
    )
    _add_footer(slide, 4)


def slide_cadastros(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Cadastros clínicos e de cardápio", "Base para padronização")
    _two_columns(
        slide,
        [
            "Pacientes, clínicas, enfermarias e leitos",
            "Dietas com categoria e grupo visual",
            "Grupos de dietas (ordem de exibição)",
            "Tipos de refeição com hora limite",
        ],
        [
            "Cardápios por dieta (popup maçã)",
            "Dieta travada no formulário do cardápio",
            "Abas: grandes, pequenas e líquidas",
            "Dietas líquidas e pratos associados",
        ],
    )
    _add_footer(slide, 5)


def slide_precos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Precificação e faturamento", "Transparência financeira")
    _bullets(
        slide,
        [
            "Grade Dieta × Tipo de Refeição com três perfis: Funcionário, Paciente e Acompanhante.",
            "Edição com opção de replicar o mesmo valor nas três colunas.",
            "Faturamento por período: espelhos, totais, fórmulas/enterais e complementares.",
            "Exportação e impressão para conferência com a diretoria e o financeiro.",
            "Mesma base do mapa — reduz divergência entre produção e cobrança.",
        ],
    )
    _add_footer(slide, 6)


def slide_estoque(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Estoque e fornecedores", "Suporte à produção")
    _bullets(
        slide,
        [
            "Produtos com código, grupo, quantidades, preços médio e último, mínimos e máximos.",
            "Unidades de medida com flags para nutrientes, UMA, estoque e pratos.",
            "Cadastro completo de fornecedores (CNPJ, contato, prazos).",
            "Tela de estoque com abas de produtos, movimentações, alertas e unidades.",
        ],
    )
    _add_footer(slide, 7)


def slide_nutrientes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Nutrientes e qualidade da informação", "Base científica operacional")
    _bullets(
        slide,
        [
            "Tabelas locais de alimentos com macros e micronutrientes (por 100 g / 100 ml).",
            "Indicadores úteis à prática: glúten, fenilalanina, coeficiente NPU, referência de consumo.",
            "Importação da FoodData Central (USDA) via ZIP/JSON — atualização estruturada do catálogo.",
            "Cardápios e substituições do mapa se apoiam nessa base cadastral.",
        ],
    )
    _add_footer(slide, 8)


def slide_rastreio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Rastreabilidade e governança", "Quem fez, quando e por quê")
    _bullets(
        slide,
        [
            "Toda alteração relevante no mapa registra usuário e data/hora.",
            "Saída do mapa exige motivo: alta médica, óbito ou transferência (com hospital destino).",
            "Histórico preservado — não se “apaga” o passado operacional.",
            "Link direto para o módulo de Auditoria da plataforma São Geraldo Service.",
            "Avisos de inconsistência (alta em mapas futuros) para decisão consciente da equipe.",
        ],
    )
    _add_footer(slide, 9)


def slide_seguranca(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Segurança e disponibilidade", "Acesso controlado e contínuo")
    _bullets(
        slide,
        [
            "Login na plataforma antes do uso dos sistemas.",
            "Serviço local em HTTP (porta 80) e HTTPS (porta 443), com certificado para ambiente interno.",
            "Inicialização simplificada: iniciar_meuapp.bat (banco + aplicação + navegador).",
            "Dados em MySQL — backup periódica do banco é prática recomendada de TI (fora da tela do módulo).",
            "Sessão identifica o profissional nas alterações do mapa.",
        ],
    )
    _add_footer(slide, 10)


def slide_beneficios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Benefícios para a Diretoria", "Controle, padronização e histórico")
    # three cards
    cards = [
        ("Controle", "Mapa único do dia por unidade; menos retrabalho e menos falha de comunicação."),
        ("Padronização", "Dietas, horários, cardápios e preços alinhados entre turnos e equipes."),
        ("Histórico", "Motivos de saída e usuário nas alterações — base para auditoria e melhoria contínua."),
    ]
    x = 0.5
    for title, body in cards:
        _add_rect(slide, Inches(x), Inches(1.9), Inches(3.9), Inches(3.6), MIST)
        _add_rect(slide, Inches(x), Inches(1.9), Inches(3.9), Inches(0.12), LIME)
        tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.3), Inches(3.4), Inches(0.5))
        r = tb.text_frame.paragraphs[0].add_run()
        _set_run(r, title, size=20, bold=True, color=BLUE)
        bb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(3.0), Inches(3.4), Inches(2.0))
        r2 = bb.text_frame.paragraphs[0].add_run()
        _set_run(r2, body, size=14, color=INK)
        x += 4.15
    _add_footer(slide, 11)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(slide, "Próximos passos (roadmap leve)", "Evolução contínua — sem métricas artificiais")
    _bullets(
        slide,
        [
            "Aprofundar movimentações e alertas de estoque no dia a dia.",
            "Consolidar painéis gerenciais (custos × mapa × faturamento) para a diretoria.",
            "Ampliar utilitários já esboçados (totalização direta, autorizações).",
            "Avaliar integração com prontuário/HIS quando a instituição priorizar.",
            "Refinar políticas de acesso e treinamento contínuo das equipes.",
        ],
    )
    _add_footer(slide, 12)


def slide_encerramento(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    _add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), LIME)

    t1 = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.6))
    r = t1.text_frame.paragraphs[0].add_run()
    _set_run(r, "Obrigado", size=36, bold=True, color=WHITE)

    t2 = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.0))
    r2 = t2.text_frame.paragraphs[0].add_run()
    _set_run(
        r2,
        "São Geraldo Service — Nutrição Hospitalar\nDocumentação completa e manual do usuário disponíveis em docs/nutricao/",
        size=16,
        color=SOFT,
    )

    t3 = slide.shapes.add_textbox(Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.6))
    r3 = t3.text_frame.paragraphs[0].add_run()
    _set_run(r3, "Perguntas e encaminhamentos", size=14, bold=True, color=LIME)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    slide_objetivo(prs)
    slide_visao(prs)
    slide_mapa(prs)
    slide_cadastros(prs)
    slide_precos(prs)
    slide_estoque(prs)
    slide_nutrientes(prs)
    slide_rastreio(prs)
    slide_seguranca(prs)
    slide_beneficios(prs)
    slide_roadmap(prs)
    slide_encerramento(prs)

    prs.save(OUT)
    print(f"OK: {OUT}")


if __name__ == "__main__":
    main()
